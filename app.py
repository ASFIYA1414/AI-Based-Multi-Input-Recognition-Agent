import streamlit as st
import pandas as pd
from sklearn.model_selection import train_test_split
from sklearn.ensemble import (
    RandomForestClassifier,
    AdaBoostClassifier,
    GradientBoostingClassifier
)
from sklearn.neighbors import KNeighborsClassifier
from sklearn.svm import SVC
from sklearn.linear_model import LogisticRegression
from sklearn.naive_bayes import GaussianNB
from sklearn.metrics import (
    confusion_matrix,
    classification_report,
    accuracy_score,
    precision_score,
    recall_score,
    f1_score,
    log_loss
)
import os
import cv2
import coll  # Ensure this is a valid module or replace with the correct import
import mediapipe as mp
import numpy as np
import time
import speech_recognition as sr
from gtts import gTTS  # Changed from pyttsx3 to gTTS for better multilingual support
import seaborn as sns
import matplotlib.pyplot as plt
import sqlite3
from googletrans import Translator
import base64
from PIL import Image
import easyocr
import logging
from io import BytesIO
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from docx import Document
from pptx import Presentation  # Added for PowerPoint file support
import PyPDF2
from streamlit_option_menu import option_menu
from moviepy.editor import VideoFileClip
import torch
from textblob import TextBlob
from transformers import pipeline, VisionEncoderDecoderModel, ViTFeatureExtractor, AutoTokenizer

# -------------------- Accessibility Reporting Setup --------------------

# Function to initialize the logging database
def init_logging_db():
    try:
        conn = sqlite3.connect('users_1_1_1.db')
        c = conn.cursor()
        # Create table for logging if it doesn't exist
        c.execute('''
            CREATE TABLE IF NOT EXISTS user_logs_1 (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                username TEXT,
                timestamp DATETIME DEFAULT CURRENT_TIMESTAMP,
                action_type TEXT,
                action_details TEXT
            )
        ''')
        conn.commit()
    except Exception as e:
        st.error(f"❌ Error initializing logging database: {e}")
        logging.error(f"Error initializing logging database: {e}")
    finally:
        conn.close()

# Function to log user actions
def log_user_action(username, action_type, action_details):
    try:
        conn = sqlite3.connect('users_1_1_1.db')
        c = conn.cursor()
        c.execute('''
            INSERT INTO user_logs_1 (username, action_type, action_details)
            VALUES (?, ?, ?)
        ''', (username, action_type, action_details))
        conn.commit()
    except Exception as e:
        st.error(f"❌ Error logging user action: {e}")
        logging.error(f"Error logging user action: {e}")
    finally:
        conn.close()

# Initialize the logging database
init_logging_db()

# -------------------- End of Accessibility Reporting Setup --------------------

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    filename='app.log'
)

# Apply custom theme
st.set_page_config(
    page_title="AI-Based Multi-Input Recognition Agent",
    page_icon=":guardsman:",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for footer and other elements
def local_css():
    st.markdown("""
    <style>
    .footer {
        position: fixed;
        left: 0;
        bottom: 0;
        width: 100%;
        background-color: #2e2e2e;
        color: white;
        text-align: center;
        padding: 10px;
    }
    .stButton>button {
        color: white;
        background-color: #4CAF50;
    }
    .stTextInput>div>div>input {
        color: black;
    }
    </style>
    """, unsafe_allow_html=True)

local_css()

# Email setup
def setup_email():
    try:
        # It's highly recommended to use environment variables or secure storage for sensitive information
        EMAIL_ADDRESS = "asfiya1404@gmail.com"  # Replace with your email
        EMAIL_PASSWORD = "sosm ltng lhgy uhdn"  # Replace with your email password or app-specific password
        s = smtplib.SMTP('smtp.gmail.com', 587)
        s.starttls()
        s.login(EMAIL_ADDRESS, EMAIL_PASSWORD)
        return s, EMAIL_ADDRESS
    except smtplib.SMTPAuthenticationError as e:
        st.error(f"SMTP Authentication Error: {e}")
        logging.error(f"SMTP Authentication Error: {e}")
        return None, None

# Function to send email
def send_email(subject, body, recipient, email_sender, email_address, username):
    try:
        if email_sender is None:
            st.error("Email sender is not set up.")
            return
        message = MIMEMultipart()
        message['From'] = email_address
        message['To'] = recipient
        message['Subject'] = subject
        message.attach(MIMEText(body, 'plain'))
        email_sender.sendmail(message['From'], message['To'], message.as_string())
        st.success("📧 Output emailed successfully.")
        # Log the email action
        log_user_action(username, "Send Email", f"Subject: {subject}, Recipient: {recipient}")
    except Exception as e:
        st.error(f"❌ Failed to send email: {e}")
        logging.error(f"Failed to send email: {e}")

# Load data and train multiple models
@st.cache_resource
def load_and_train_models(username=None):
    try:
        local_path = os.path.dirname(os.path.realpath('__file__'))
        file_name = 'data1.csv'  # Ensure this file is in the same directory
        data_path = os.path.join(local_path, file_name)
        df = pd.read_csv(data_path)

        units_in_data = 28  # Number of units in data
        titles = [f"unit-{i}" for i in range(units_in_data)]
        X = df[titles]
        y = df['letter']

        # Split data
        X_train, X_test, y_train, y_test = train_test_split(
            X, y, test_size=0.3, random_state=42
        )

        # Define classifiers
        classifiers = {
            'Random Forest': RandomForestClassifier(random_state=42),
            'K-Nearest Neighbors': KNeighborsClassifier(),
            'Support Vector Machine': SVC(probability=True),
            'Logistic Regression': LogisticRegression(max_iter=1000),
            'Naive Bayes': GaussianNB(),
            'AdaBoost': AdaBoostClassifier(random_state=42),
            'Gradient Boosting': GradientBoostingClassifier(random_state=42)
        }

        # For storing results
        metrics = {}
        confusion_matrices = {}
        classification_reports = {}
        trained_classifiers = {}

        for name, clf in classifiers.items():
            with st.spinner(f"📈 Training {name}..."):
                start_time = time.time()
                clf.fit(X_train, y_train)
                training_time = time.time() - start_time

                y_pred = clf.predict(X_test)
                accuracy = accuracy_score(y_test, y_pred)
                precision = precision_score(
                    y_test, y_pred, average='weighted', zero_division=0
                )
                recall = recall_score(
                    y_test, y_pred, average='weighted', zero_division=0
                )
                f1 = f1_score(
                    y_test, y_pred, average='weighted', zero_division=0
                )
                # Compute Log Loss if possible
                try:
                    y_pred_proba = clf.predict_proba(X_test)
                    logloss = log_loss(y_test, y_pred_proba)
                except:
                    logloss = None  # Some classifiers don't support predict_proba

                cm = confusion_matrix(y_test, y_pred)
                report = classification_report(
                    y_test, y_pred, output_dict=True
                )
                df_report = pd.DataFrame(report).transpose()

                metrics[name] = {
                    'Accuracy': accuracy,
                    'Precision': precision,
                    'Recall': recall,
                    'F1 Score': f1,
                    'Log Loss': logloss,
                    'Training Time': training_time
                }
                confusion_matrices[name] = cm
                classification_reports[name] = df_report
                trained_classifiers[name] = clf

        st.success("🎉 All models trained successfully!")
        # Log the model training action
        log_user_action(username, "Train Models", f"Trained classifiers: {list(classifiers.keys())}")
        return trained_classifiers, metrics, confusion_matrices, classification_reports

    except Exception as e:
        st.error(f"❌ Error in loading and training models: {e}")
        logging.error(f"Error in loading and training models: {e}")

# Database interaction functions without password hashing
def create_user(username, password):
    try:
        conn = sqlite3.connect('users_1_1_1.db')
        c = conn.cursor()
        c.execute('CREATE TABLE IF NOT EXISTS users_1_1_1 (username TEXT PRIMARY KEY, password TEXT)')
        c.execute('INSERT INTO users_1_1_1 (username, password) VALUES (?, ?)', (username, password))
        conn.commit()
        st.success("✅ User created successfully.")
        # Log the user creation
        log_user_action(username, "Create User", "User account created.")
    except sqlite3.IntegrityError:
        st.error("⚠️ Username already exists.")
    except sqlite3.OperationalError as e:
        st.error(f"❌ Database error: {e}")
        logging.error(f"Database error: {e}")
    finally:
        conn.close()

def authenticate_user(username, password):
    try:
        conn = sqlite3.connect('users_1_1_1.db')
        c = conn.cursor()
        c.execute('SELECT password FROM users_1_1_1 WHERE username = ?', (username,))
        result = c.fetchone()
        if result and password == result[0]:
            # Log the successful login
            log_user_action(username, "Login", "User logged in successfully.")
            return True
        else:
            # Log the failed login attempt
            log_user_action(username, "Login Failed", "Invalid username or password.")
            return False
    except sqlite3.OperationalError as e:
        st.error(f"❌ Database error: {e}")
        logging.error(f"Database error: {e}")
        return False
    finally:
        conn.close()

# Hand sign prediction setup
mp_drawing = mp.solutions.drawing_utils
mp_hands = mp.solutions.hands

# Function to get prediction from hand landmarks
def get_prediction(image, clf):
    with mp_hands.Hands(
        min_detection_confidence=0.5,
        min_tracking_confidence=0.5
    ) as hands:
        ImageData = coll.ImageToDistanceData(image, hands)
        DistanceData = ImageData['Distance-Data']
        if len(DistanceData) == 0:
            st.write("❌ Error: DistanceData is empty.")
            return "UNKNOWN"
        prediction = clf.predict([DistanceData])
        return prediction[0]

# Function to translate text with error handling
def translate_text(text, target_language, username, action_type="Translate Text"):
    translator = Translator()
    try:
        translated = translator.translate(text, dest=target_language)
        # Log the translation action
        log_user_action(username, action_type, f"Original Text: {text} | Translated Text: {translated.text}")
        return translated.text
    except Exception as e:
        st.write(f"❌ Error in translation: {e}")
        logging.error(f"Error in translation: {e}")
        return text

# Function to handle hand signs with translation option
def handle_hand_signs(SpelledWord, target_language, email_sender, email_address, username):
    predefined_texts = {
        "HAVE A GREAT DAY!": "HAVE A GREAT DAY!",
        "EMERGENCY": "EMERGENCY",
        "A CUP OF COFFEE?": "A CUP OF COFFEE?",
        "HAVE A BREAKFAST!": "HAVE A BREAKFAST!",
        "GOOD NIGHT! SLEEP WELL.": "GOOD NIGHT! SLEEP WELL.",
        "HOW ARE YOU?": "HOW ARE YOU?"
    }

    detected_text = predefined_texts.get(SpelledWord, "GIVE THE HAND SIGNS")
    st.write(f"**Detected Hand Sign:** {detected_text}")
    translated_text = translate_text(detected_text, target_language, username)
    st.write(f"**Translated Text:** {translated_text}")

    # Convert text to speech
    text_to_voice(translated_text, target_language)

    # Perform sentiment analysis
    sentiment_analysis(translated_text, username)

    # Log the recognized gesture
    log_user_action(username, "Recognized Gesture", f"Gesture: {SpelledWord} | Translated: {translated_text}")

    # Option to send email
    if st.button("Send Output via Email"):
        subject = "Hand Sign Recognition Output"
        body = f"Detected Hand Sign: {detected_text}\nTranslated Text: {translated_text}"
        recipient = email_address  # Sending to the sender's email
        send_email(subject, body, recipient, email_sender, email_address, username)

# Function to handle audio input from microphone with noise reduction
def audio_to_text_from_microphone(target_language, email_sender, elmail_address, username):
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        st.write("Adjusting for ambient noise, please wait...")
        recognizer.adjust_for_ambient_noise(source, duration=1)
        st.write("Listening for live audio input...")
        audio = recognizer.listen(source)
        try:
            with st.spinner("Recognizing speech..."):
                text = recognizer.recognize_google(audio)
            st.write(f"**You said:** {text}")
            translated_text = translate_text(text, target_language, username)
            st.write(f"**Translated Text:** {translated_text}")
            # Convert text to speech
            text_to_voice(translated_text, target_language)
            # Perform sentiment analysis
            sentiment_analysis(translated_text, username)
            # Log the audio translation
            log_user_action(username, "Audio Translation", f"Original Text: {text} | Translated Text: {translated_text}")
            return text, translated_text
        except sr.UnknownValueError:
            st.write("Speech Recognition could not understand the audio.")
            logging.error("Speech Recognition could not understand the audio.")
            # Log the failed audio translation
            log_user_action(username, "Audio Translation Failed", "Could not understand audio.")
            return None, None
        except sr.RequestError as e:
            st.write(f"⚠️ Could not request results; {e}")
            logging.error(f"Speech Recognition request error: {e}")
            # Log the request error
            log_user_action(username, "Audio Translation Error", f"Request error: {e}")
            return None, None

# Function to convert text to voice using gTTS
def text_to_voice(text, language_code):
    try:
        tts = gTTS(text=text, lang=language_code, slow=False)
        temp_audio = "temp_audio.mp3"
        tts.save(temp_audio)
        with open(temp_audio, "rb") as audio_file:
            audio_bytes = audio_file.read()
            st.audio(audio_bytes, format="audio/mp3")
        os.remove(temp_audio)
    except Exception as e:
        st.write(f"⚠️ Error in text-to-voice conversion: {e}")
        logging.error(f"Error in text-to-voice conversion: {e}")

# Expanded list of languages (Google Translate supports around 100 languages)
language_options = {
    "Afrikaans": "af",
    "Albanian": "sq",
    "Amharic": "am",
    "Arabic": "ar",
    "Armenian": "hy",
    "Azerbaijani": "az",
    "Basque": "eu",
    "Belarusian": "be",
    "Bengali": "bn",
    "Bosnian": "bs",
    "Bulgarian": "bg",
    "Catalan": "ca",
    "Cebuano": "ceb",
    "Chinese (Simplified)": "zh-CN",
    "Chinese (Traditional)": "zh-TW",
    "Corsican": "co",
    "Croatian": "hr",
    "Czech": "cs",
    "Danish": "da",
    "Dutch": "nl",
    "English": "en",
    "Esperanto": "eo",
    "Estonian": "et",
    "Finnish": "fi",
    "French": "fr",
    "Galician": "gl",
    "Georgian": "ka",
    "German": "de",
    "Greek": "el",
    "Gujarati": "gu",
    "Haitian Creole": "ht",
    "Hausa": "ha",
    "Hebrew": "iw",
    "Hindi": "hi",
    "Hmong": "hmn",
    "Hungarian": "hu",
    "Icelandic": "is",
    "Igbo": "ig",
    "Indonesian": "id",
    "Irish": "ga",
    "Italian": "it",
    "Japanese": "ja",
    "Javanese": "jw",
    "Kannada": "kn",
    "Kazakh": "kk",
    "Khmer": "km",
    "Korean": "ko",
    "Kurdish (Kurmanji)": "ku",
    "Kyrgyz": "ky",
    "Lao": "lo",
    "Latin": "la",
    "Latvian": "lv",
    "Lithuanian": "lt",
    "Luxembourgish": "lb",
    "Macedonian": "mk",
    "Malagasy": "mg",
    "Malay": "ms",
    "Malayalam": "ml",
    "Maltese": "mt",
    "Maori": "mi",
    "Marathi": "mr",
    "Mongolian": "mn",
    "Nepali": "ne",
    "Norwegian": "no",
    "Pashto": "ps",
    "Persian": "fa",
    "Polish": "pl",
    "Portuguese": "pt",
    "Punjabi": "pa",
    "Romanian": "ro",
    "Russian": "ru",
    "Serbian": "sr",
    "Slovak": "sk",
    "Slovenian": "sl",
    "Somali": "so",
    "Spanish": "es",
    "Swahili": "sw",
    "Swedish": "sv",
    "Tamil": "ta",
    "Telugu": "te",
    "Thai": "th",
    "Turkish": "tr",
    "Ukrainian": "uk",
    "Urdu": "ur",
    "Vietnamese": "vi",
    "Welsh": "cy",
    "Xhosa": "xh",
    "Yiddish": "yi",
    "Yoruba": "yo",
    "Zulu": "zu"
}

# Function to load the image captioning model
@st.cache_resource
def load_image_captioning_model():
    try:
        model = VisionEncoderDecoderModel.from_pretrained("nlpconnect/vit-gpt2-image-captioning")
        feature_extractor = ViTFeatureExtractor.from_pretrained("nlpconnect/vit-gpt2-image-captioning")
        tokenizer = AutoTokenizer.from_pretrained("nlpconnect/vit-gpt2-image-captioning")
        device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        model.to(device)
        return model, feature_extractor, tokenizer, device
    except Exception as e:
        st.error(f"❌ Error loading image captioning model: {e}")
        logging.error(f"Error loading image captioning model: {e}")
        return None, None, None, None

# Function to generate image caption
def generate_image_caption(image, model, feature_extractor, tokenizer, device, username, max_length=16, num_beams=4):
    try:
        if model is None:
            st.error("Image captioning model is not loaded.")
            return "No Caption Available"

        # Preprocess the image
        image = image.convert("RGB")
        pixel_values = feature_extractor(images=image, return_tensors="pt").pixel_values
        pixel_values = pixel_values.to(device)

        # Generate caption
        with torch.no_grad():
            output_ids = model.generate(pixel_values, max_length=max_length, num_beams=num_beams)
        caption = tokenizer.decode(output_ids[0], skip_special_tokens=True).strip()

        # Log the image captioning action
        log_user_action(username, "Image Captioning", f"Caption: {caption}")

        return caption
    except Exception as e:
        st.error(f"❌ Error generating image caption: {e}")
        logging.error(f"Error generating image caption: {e}")
        return "No Caption Available"

# Function to load the sentiment analysis pipeline
@st.cache_resource
def load_sentiment_pipeline():
    try:
        # Load a multilingual sentiment analysis model using PyTorch
        sentiment_pipeline = pipeline(
            "sentiment-analysis",
            model="nlptown/bert-base-multilingual-uncased-sentiment",
            framework='pt'  # Force using PyTorch
        )
        return sentiment_pipeline
    except Exception as e:
        st.error(f"❌ Error loading sentiment analysis model: {e}")
        logging.error(f"Error loading sentiment analysis model: {e}")
        return None

# Function to perform sentiment analysis
def sentiment_analysis(text, username):
    try:
        sentiment_pipeline = load_sentiment_pipeline()
        if sentiment_pipeline is None:
            st.write("Sentiment analysis model is not available.")
            return

        # Get sentiment results
        results = sentiment_pipeline(text)

        # Process results
        sentiment_mapping = {
            '1 star': "Very Negative",
            '2 stars': "Negative",
            '3 stars': "Neutral",
            '4 stars': "Positive",
            '5 stars': "Very Positive"
        }

        emoji_mapping = {
            "Very Negative": "😠",
            "Negative": "😞",
            "Neutral": "😐",
            "Positive": "🙂",
            "Very Positive": "😃"
        }

        for result in results:
            label = result['label']
            score = result['score']
            sentiment = sentiment_mapping.get(label, "Unknown")
            emoji = emoji_mapping.get(sentiment, "")
            st.write(f"**Sentiment Analysis:**")
            st.write(f"**Sentiment:** {sentiment} {emoji}")
            st.write(f"**Confidence Score:** {score:.2f}")

            # Log the sentiment analysis action
            log_user_action(username, "Sentiment Analysis", f"Sentiment: {sentiment}, Score: {score}")

            # Optionally, visualize the sentiment
            fig, ax = plt.subplots(figsize=(8, 4))  # Adjusted figure size
            sentiments = list(sentiment_mapping.values())
            scores = [0, 0, 0, 0, 0]
            # Extract the number of stars from the label using seaborn we designed sentiments confidence
            if 'stars' in label or 'star' in label:
                star_number = int(label.split()[0])
                if 1 <= star_number <= 5:
                    scores[star_number - 1] = score
            sns.barplot(x=sentiments, y=scores, ax=ax, palette="viridis")
            ax.set_ylim(0, 1)
            ax.set_ylabel("Confidence Score")
            ax.set_title("Sentiment Confidence Levels")
            plt.tight_layout()
            st.pyplot(fig, use_container_width=True)

    except Exception as e:
        st.write(f"⚠️ Error in sentiment analysis: {e}")
        logging.error(f"Error in sentiment analysis: {e}")

# Main Streamlit app logic
def main():
    # Navbar using streamlit_option_menu
    with st.sidebar:
        selected = option_menu(
            menu_title="Main Menu",  # Direct English text
            options=["Home", "Authentication", "Tasks", "Reports", "About"],
            icons=["house", "person", "gear", "bar-chart", "info-circle"],
            menu_icon="cast",
            default_index=0,
            styles={
                "container": {"padding": "0!important", "background-color": "#2e2e2e"},
                "icon": {"color": "white", "font-size": "25px"},
                "nav-link": {"font-size": "16px", "text-align": "left", "margin": "0px", "--hover-color": "#555"},
                "nav-link-selected": {"background-color": "#1eff96e1"},
            }
        )

    if selected == "Home":
        st.title("Welcome to AI-Based Multi-Input Recognition Agent")
        st.markdown("""
            This platform allows you to:
            - 🤚 **Recognize and translate hand signs.**
            - 🎧 **Process audio inputs** from microphone or files.
            - 🪄 **Extract and translate text** from images, videos, and live video.
            - 📄 **Translate text and documents** into over 100 languages.
            - 🖋️ **Generate captions for images** to understand and translate visual content effortlessly.
            - 💡 **Perform sentiment analysis** on extracted or input text.
            - 📧 **Receive all outputs via email.**
        """)
        # Add a header image (replace the URL with your own image or upload an image)
        st.image("https://images.unsplash.com/photo-1557683316-973673baf926", use_column_width=True)

    elif selected == "Authentication":
        if 'user_authenticated' not in st.session_state:
            st.session_state.user_authenticated = False
            st.session_state.previous_page = None

        if not st.session_state.user_authenticated:
            auth_menu = option_menu(
                menu_title="Authentication",  # Direct English text
                options=["Sign Up", "Login"],
                icons=["person-plus", "box-arrow-in-right"],
                menu_icon="key",
                default_index=0,
                orientation="vertical",
                styles={
                    "container": {"padding": "0!important", "background-color": "#2e2e2e"},
                    "icon": {"color": "orange", "font-size": "20px"},
                    "nav-link": {"font-size": "14px", "text-align": "left", "margin": "0px", "--hover-color": "#555"},
                    "nav-link-selected": {"background-color": "#1effa9be"},
                }
            )

            if auth_menu == "Sign Up":
                st.subheader("Create New Account")
                new_username = st.text_input("Username")
                new_password = st.text_input("Password", type='password')
                if st.button("Sign Up"):
                    if new_username and new_password:
                        create_user(new_username, new_password)
                    else:
                        st.warning("Please enter both username and password.")

            elif auth_menu == "Login":
                st.subheader("Login")
                username = st.text_input("Username")
                password = st.text_input("Password", type='password')
                if st.button("Login"):
                    if authenticate_user(username, password):
                        st.session_state.user_authenticated = True
                        st.session_state.username = username
                        st.success("Logged in successfully!")
                    else:
                        st.error("Invalid username or password.")
        else:
            st.success(f"👤 Logged in as {st.session_state.username}")
            if st.button("Logout"):
                st.session_state.user_authenticated = False
                st.session_state.username = None
                st.rerun()


    elif selected == "Tasks":
        if 'user_authenticated' not in st.session_state or not st.session_state.user_authenticated:
            st.warning("Please login to access the tasks.")
        else:
            username = st.session_state.username  # Retrieve the logged-in username
            # Sidebar for language and task selection
            with st.sidebar:
                language = st.selectbox("Select Language", list(language_options.keys()))
                target_language = language_options[language]

                task = st.selectbox("Choose Task", [
                    "Hand Sign Recognition",
                    "Audio Translation",
                    "Image and Video Translation",
                    "Text and Document Translation",
                    "Image Captioning"  # New Task Added
                ])

                # Email setup
                email_sender, email_address = setup_email()

            # Load and train models
            trained_classifiers, metrics, confusion_matrices, classification_reports = load_and_train_models(username)

            # Display model evaluation in an expander for Hand Sign Recognition
            if task == "Hand Sign Recognition":
                with st.expander("Model Evaluation Metrics"):
                    st.write("Model Comparison")

                    # Create a DataFrame for metrics
                    metrics_df = pd.DataFrame(metrics).transpose()
                    st.dataframe(metrics_df.style.highlight_max(axis=0))

                    # Plot bar charts for each metric
                    metric_options = ['Accuracy', 'Precision', 'Recall', 'F1 Score', 'Training Time']
                    for metric in metric_options:
                        st.write(f"### 📉 {metric}")
                        fig, ax = plt.subplots(figsize=(8, 4))  # Adjusted figure size
                        sns.barplot(x=metrics_df.index, y=metrics_df[metric], ax=ax, palette="viridis")
                        ax.set_ylabel(metric)
                        ax.set_xlabel("Classifier")
                        plt.xticks(rotation=45, ha='right')  # Improved label alignment
                        plt.tight_layout()
                        st.pyplot(fig, use_container_width=True)

                    # Display confusion matrix for each classifier
                    for clf_name in confusion_matrices:
                        st.write(f"#### 🔍 {clf_name} Confusion Matrix:")
                        cm = confusion_matrices[clf_name]
                        fig_cm, ax_cm = plt.subplots(figsize=(8, 6))  # Reduced figure size
                        sns.heatmap(cm, annot=True, fmt=".0f", linewidths=.5, square=True, cmap='Blues', ax=ax_cm)
                        ax_cm.set_title(f'{clf_name} Confusion Matrix')
                        ax_cm.set_xlabel("Predicted")
                        ax_cm.set_ylabel("Actual")
                        plt.tight_layout()
                        st.pyplot(fig_cm, use_container_width=True)

            # Load the image captioning model once
            image_caption_model, feature_extractor, tokenizer, device = load_image_captioning_model()

            st.header(f"### 🚀 {task}")
            # Removed: ui_lang parameter

            if task == "Hand Sign Recognition":
                if st.button("Start Hand Sign Recognition"):
                    st.info("Starting webcam for hand sign recognition. Please allow access to your camera.")
                    # Select the best performing classifier based on accuracy
                    best_clf_name = max(metrics, key=lambda x: metrics[x]['Accuracy'])
                    clf = trained_classifiers[best_clf_name]
                    st.write(f"Using **{best_clf_name}** for predictions.")
                    hand_sign_to_text(target_language, clf, email_sender, email_address, username)

            elif task == "Audio Translation":
                audio_option = st.radio("Select Audio Input Type", ["Live Audio Input", "Upload Audio File"])
                if audio_option == "Live Audio Input":
                    if st.button("Start Live Audio Input"):
                        text, translated_text = audio_to_text_from_microphone(
                            target_language, email_sender, email_address, username
                        )
                        if text and translated_text:
                            st.balloons()
                            if st.button("Send Output via Email"):
                                subject = "Live Audio Input Output"
                                body = f"Original Text: {text}\nTranslated Text: {translated_text}"
                                recipient = email_address  # Sending to the sender's email
                                send_email(subject, body, recipient, email_sender, email_address, username)
                elif audio_option == "Upload Audio File":
                    uploaded_file = st.file_uploader("Upload an audio file", type=["wav", "mp3", "flac"])
                    if uploaded_file is not None:
                        with st.spinner("Processing the uploaded audio file..."):
                            file_path = os.path.join('temp', uploaded_file.name)
                            if not os.path.exists('temp'):
                                os.makedirs('temp')
                            with open(file_path, "wb") as f:
                                f.write(uploaded_file.read())
                            text, translated_text = audio_to_text_from_file(
                                file_path, target_language, email_sender, email_address, username
                            )
                        if text and translated_text:
                            st.balloons()
                            if st.button("Send Output via Email"):
                                subject = "Audio File Processing Output"
                                body = f"Extracted Text: {text}\nTranslated Text: {translated_text}"
                                recipient = email_address  # Sending to the sender's email
                                send_email(subject, body, recipient, email_sender, email_address, username)

            elif task == "Image and Video Translation":
                media_option = st.radio("Select Media Type", ["Image", "Video File", "Live Video"])
                if media_option == "Image":
                    uploaded_image = st.file_uploader("Upload an image file", type=["jpg", "jpeg", "png"])
                    if uploaded_image is not None:
                        image = Image.open(uploaded_image)
                        st.image(image, caption='🖼️ Uploaded Image', use_column_width=True)
                        # Extract text from image
                        text, translated_text = image_to_text(
                            image, target_language, email_sender, email_address, username
                        )
                        if text.strip():
                            # If text is found, process text translation
                            st.write(f"**Text extracted from image:**")
                            st.write(text)
                            st.write(f"**Translated Text:**")
                            st.write(translated_text)
                            # Convert text to speech
                            text_to_voice(translated_text, target_language)
                            # Perform sentiment analysis
                            sentiment_analysis(translated_text, username)
                            # Provide download option
                            st.download_button(
                                label="Download Translated Text",
                                data=translated_text,
                                file_name='translated_text.txt',
                                mime='text/plain',
                                key="download_image_translation"
                            )
                            # Option to send email with text
                            if st.button("Send Output via Email"):
                                subject = "Image Text Translation Output"
                                body = f"**Extracted Text:** {translated_text}"
                                recipient = email_address  # Sending to the sender's email
                                send_email(subject, body, recipient, email_sender, email_address, username)
                        else:
                            # If no text is found, perform image captioning
                            with st.spinner("Generating image caption..."):
                                caption = generate_image_caption(image, image_caption_model, feature_extractor, tokenizer, device, username)
                            st.write(f"**Generated Caption:** {caption}")
                            # Translate caption
                            translated_caption = translate_text(caption, target_language, username, action_type="Image Captioning")
                            st.write(f"**Translated Caption:** {translated_caption}")
                            # Convert caption to speech
                            text_to_voice(translated_caption, target_language)
                            # Perform sentiment analysis on caption
                            sentiment_analysis(translated_caption, username)
                            # Provide download option for caption
                            st.download_button(
                                label="Download Translated Caption",
                                data=translated_caption,
                                file_name='translated_caption.txt',
                                mime='text/plain'
                            )
                            # Option to send email with caption
                            if st.button("Send Caption via Email"):
                                subject = "Image Captioning Output"
                                body = f"**Generated Caption:** {translated_caption}"
                                recipient = email_address  # Sending to the sender's email
                                send_email(subject, body, recipient, email_sender, email_address, username)

                elif media_option == "Video File":
                    uploaded_video = st.file_uploader("Upload a video file", type=["mp4", "avi", "mov", "mkv"])
                    if uploaded_video is not None:
                        with st.spinner("Processing the uploaded video file..."):
                            text, translated_text = video_file_to_text(
                                uploaded_video, target_language, email_sender, email_address, username
                            )
                        if text and translated_text:
                            st.balloons()
                            if st.button("Send Output via Email"):
                                subject = "Video File Processing Output"
                                body = f"Extracted Text: {text}\nTranslated Text: {translated_text}"
                                recipient = email_address  # Sending to the sender's email
                                send_email(subject, body, recipient, email_sender, email_address, username)
                elif media_option == "Live Video":
                    if st.button("Start Live Video Translation"):
                        text, translated_text = live_video_translation(
                            target_language, email_sender, email_address, username
                        )
                        if text and translated_text:
                            st.balloons()
                            if st.button("Send Output via Email"):
                                subject = "Live Video Translation Output"
                                body = f"Captured Audio Text: {text}\nTranslated Text: {translated_text}"
                                recipient = email_address  # Sending to the sender's email
                                send_email(subject, body, recipient, email_sender, email_address, username)

            elif task == "Text and Document Translation":
                doc_option = st.radio("Select Input Type", ["Enter Text", "Upload Document"])
                if doc_option == "Enter Text":
                    user_text = st.text_area("Enter text for translation")
                    if st.button("Translate Text"):
                        if user_text:
                            translated_text = translate_text(user_text, target_language, username)
                            st.write(f"**Translated Text:** {translated_text}")
                            # Convert text to speech
                            text_to_voice(translated_text, target_language)
                            # Perform sentiment analysis
                            sentiment_analysis(translated_text, username)
                            # Provide download option
                            st.download_button(
                                label="Download Translated Text",
                                data=translated_text,
                                file_name='translated_text.txt',
                                mime='text/plain'
                            )
                            # Option to send email
                            if st.button("Send Output via Email"):
                                subject = "Text Translation Output"
                                body = f"Original Text: {user_text}\nTranslated Text: {translated_text}"
                                recipient = email_address  # Sending to the sender's email
                                send_email(subject, body, recipient, email_sender, email_address, username)
                        else:
                            st.warning("Please enter some text.")
                elif doc_option == "Upload Document":
                    doc_type = st.selectbox("Select Document Type", ["Text File", "Word Document", "PDF File", "PowerPoint Presentation"])
                    if doc_type == "Text File":
                        uploaded_text_file = st.file_uploader("Upload a text file", type=["txt"])
                        if uploaded_text_file is not None:
                            text, translated_text = text_file_to_text(  
                                uploaded_text_file, target_language, email_sender, email_address, username
                            )
                            if text and translated_text:
                                if st.button("Send Output via Email"):
                                    subject = "Text File Translation Output"
                                    body = f"Original Text: {text}\nTranslated Text: {translated_text}"
                                    recipient = email_address  # Sending to the sender's email
                                    send_email(subject, body, recipient, email_sender, email_address, username)
                    elif doc_type == "Word Document":
                        uploaded_word_file = st.file_uploader("Upload a Word document", type=["docx"])
                        if uploaded_word_file is not None:
                            text, translated_text = word_file_to_text(
                                uploaded_word_file, target_language, email_sender, email_address, username
                            )
                            if text and translated_text:
                                if st.button("Send Output via Email"):
                                    subject = "Word Document Translation Output"
                                    body = f"Original Text: {text}\nTranslated Text: {translated_text}"
                                    recipient = email_address  # Sending to the sender's email
                                    send_email(subject, body, recipient, email_sender, email_address, username)
                    elif doc_type == "PDF File":
                        uploaded_pdf_file = st.file_uploader("Upload a PDF file", type=["pdf"])
                        if uploaded_pdf_file is not None:
                            text, translated_text = pdf_file_to_text(
                                uploaded_pdf_file, target_language, email_sender, email_address, username
                            )
                            if text and translated_text:
                                if st.button("Send Output via Email"):
                                    subject = "PDF File Translation Output"
                                    body = f"Original Text: {text}\nTranslated Text: {translated_text}"
                                    recipient = email_address  # Sending to the sender's email
                                    send_email(subject, body, recipient, email_sender, email_address, username)
                    elif doc_type == "PowerPoint Presentation":
                        uploaded_ppt_file = st.file_uploader("Upload a PowerPoint presentation", type=["pptx"])
                        if uploaded_ppt_file is not None:
                            text, translated_text = ppt_file_to_text(
                                uploaded_ppt_file, target_language, email_sender, email_address, username
                            )
                            if text and translated_text:
                                if st.button("Send Output via Email"):
                                    subject = "PowerPoint Presentation Translation Output"
                                    body = f"Original Text: {text}\nTranslated Text: {translated_text}"
                                    recipient = email_address  # Sending to the sender's email
                                    send_email(subject, body, recipient, email_sender, email_address, username)

            elif task == "Image Captioning":
                st.subheader("Image Captioning")
                uploaded_image = st.file_uploader("Upload an image for captioning", type=["jpg", "jpeg", "png"])
                if uploaded_image is not None:
                    image = Image.open(uploaded_image)
                    st.image(image, caption='🖼️ Uploaded Image', use_column_width=True)
                    with st.spinner("Generating caption..."):
                        caption = generate_image_caption(image, image_caption_model, feature_extractor, tokenizer, device, username)
                    st.success("Caption generated successfully!")
                    st.write(f"**Generated Caption:** {caption}")

                    # Translate the caption
                    translated_caption = translate_text(caption, target_language, username, action_type="Image Captioning")
                    st.write(f"**Translated Caption:** {translated_caption}")

                    # Convert translated text to speech
                    text_to_voice(translated_caption, target_language)

                    # Perform sentiment analysis
                    sentiment_analysis(translated_caption, username)

                    # Provide download option for the caption
                    st.download_button(
                        label="Download Caption",
                        data=translated_caption,
                        file_name='image_caption.txt',
                        mime='text/plain'
                    )

                    # Option to send caption via email
                    if st.button("Send Caption via Email"):
                        subject = "Image Captioning Output"
                        body = f"**Generated Caption:** {caption}\n**Translated Caption:** {translated_caption}"
                        recipient = email_address  # Sending to the sender's email
                        send_email(subject, body, recipient, email_sender, email_address, username)

    elif selected == "Reports":
        if 'user_authenticated' not in st.session_state or not st.session_state.user_authenticated:
            st.warning("Please login to access the reports.")
        else:
            username = st.session_state.username
            st.title("📊 Accessibility Reports")

            # Connect to the database
            try:
                conn = sqlite3.connect('users_1_1_1.db')
                c = conn.cursor()

                # Fetch logs for the current user
                c.execute('''
                    SELECT timestamp, action_type, action_details
                    FROM user_logs_1
                    WHERE username = ?
                    ORDER BY timestamp DESC
                ''', (username,))
                logs = c.fetchall()

                # Convert to DataFrame
                logs_df = pd.DataFrame(logs, columns=['Timestamp', 'Action Type', 'Action Details'])

                if logs_df.empty:
                    st.info("No logs available to display.")
                else:
                    st.dataframe(logs_df)

                    # Filter options
                    st.sidebar.header("Filter Reports")
                    action_type_filter = st.sidebar.multiselect(
                        "Select Action Types",
                        options=logs_df['Action Type'].unique(),
                        default=logs_df['Action Type'].unique()
                    )

                    date_filter = st.sidebar.date_input(
                        "Select Date Range",
                        []
                    )

                    if date_filter:
                        start_date = date_filter[0]
                        end_date = date_filter[1]
                        logs_df['Timestamp'] = pd.to_datetime(logs_df['Timestamp'])
                        filtered_df = logs_df[
                            (logs_df['Timestamp'] >= pd.to_datetime(start_date)) &
                            (logs_df['Timestamp'] <= pd.to_datetime(end_date))
                        ]
                    else:
                        filtered_df = logs_df

                    if action_type_filter:
                        filtered_df = filtered_df[filtered_df['Action Type'].isin(action_type_filter)]

                    st.write("### Filtered Logs")
                    st.dataframe(filtered_df)

                    # Download option
                    csv = filtered_df.to_csv(index=False).encode('utf-8')
                    st.download_button(
                        label="Download Report as CSV",
                        data=csv,
                        file_name='accessibility_report.csv',
                        mime='text/csv'
                    )

            except Exception as e:
                st.error(f"❌ Error fetching reports: {e}")
                logging.error(f"Error fetching reports: {e}")
            finally:
                conn.close()

    elif selected == "About":
        st.title("About")
        st.markdown("""
            ### AI-Based Multi-Input Recognition Agent

            **Overview:**
            This application integrates various AI and machine learning technologies to provide a comprehensive multi-input recognition agent.

            **Features:**
            - **Hand Sign Recognition:** Translate hand gestures into text and speech using computer vision techniques.
            - **Audio Translation:** Convert live or recorded audio into text and translate it into over 100 languages.
            - **Image and Video Translation:** Extract text from images, videos, and live video streams, then translate and analyze it.
            - **Text and Document Translation:** Translate text from various document formats, including plain text, Word documents, PDFs, and PowerPoint presentations.
            - **Image Captioning:** Automatically generate descriptive captions for uploaded images.
            - **Sentiment Analysis:** Analyze the sentiment of the input text to understand emotions and opinions.
            - **Email Integration:** Receive all outputs via email for record-keeping and further use.
            - **Accessibility Reporting:** Generate reports on how the app assists users_1_1_1 by logging translations, recognized gestures, and sentiment analysis results.

            **Technologies Used:**
            - **Streamlit:** For building the web application interface.
            - **Scikit-learn:** For machine learning models and evaluation.
            - **Mediapipe & OpenCV:** For hand sign recognition and image processing.
            - **SpeechRecognition & gTTS:** For audio processing and text-to-speech functionality.
            - **Pytesseract:** For OCR (Optical Character Recognition) to extract text from images.
            - **Googletrans:** For translation services across multiple languages.
            - **TextBlob:** For sentiment analysis to gauge the emotional tone.
            - **Transformers (Hugging Face):** For image captioning and sentiment analysis models.
            - **PyTorch:** Backend for model inference.

            **Definitions:**
            - **OCR (Optical Character Recognition):** Technology that recognizes text within a digital image.
            - **Sentiment Analysis:** Computational study of opinions, sentiments, and emotions expressed in text.
            - **Hand Sign Recognition:** The process of interpreting human hand gestures via mathematical algorithms.
            - **Image Captioning:** Automatically generating descriptive text for images using AI models.

            **Authors:**
            | Student ID | Name                             | Department | Specialization |
            |------------|----------------------------------|------------|-----------------|
            | 22BCE20226 | Nischal Kota                     | CSE        | Core            |
            | 22BCE9313  | Shubhra Yadav                    | CSE        | Core            |
            | 22BCE9699  | Asfiya Madina Shaik              | CSE        | Core            |
            | 22BCE20192 | Harshith Sallangi                | CSE        | Core            |

            **Guide:**
            - **Dr. Deepan RamKumar P**

            **Contact:**
            For any queries or support, please contact [asfiya1404@gmail.com][shubhrayadav05@gmail.com].
        """)
        # Add an about image (replace the URL with your own image or upload an image)
        st.image("https://images.unsplash.com/photo-1521791136064-7986c2920216", use_column_width=True)

    # Footer
    footer = """
    <style>
    .footer {
        position: fixed;
        left: 0;
        bottom: 0;
        width: 100%;
        background-color: #2e2e2e;
        color: white;
        text-align: center;
        padding: 10px;
    }
    </style>
    <div class="footer">
        <p>© 2025 AI-Based Multi-Input Recognition Agent. All rights reserved.</p>
    </div>
    """
    st.markdown(footer, unsafe_allow_html=True)

# Function to handle hand sign recognition
def hand_sign_to_text(target_language, clf, email_sender, email_address, username):
    cap = cv2.VideoCapture(0)
    SpelledWord = ""
    stframe = st.empty()
    progress_bar = st.progress(0)
    for percent_complete in range(100):
        time.sleep(0.01)
        progress_bar.progress(percent_complete + 1)
    while cap.isOpened():
        success, image = cap.read()
        if not success:
            st.write("Ignoring empty camera frame.")
            continue

        try:
            SpelledWord = get_prediction(image, clf)
            # Handle the recognized hand sign
            handle_hand_signs(SpelledWord, target_language, email_sender, email_address, username)

            # Overlay the prediction on the image
            cv2.putText(
                image, SpelledWord, (50, 50),
                cv2.FONT_HERSHEY_SIMPLEX, 1, (124, 252, 0), 2, cv2.LINE_AA
            )

            # Display the frame in Streamlit
            stframe.image(
                cv2.cvtColor(image, cv2.COLOR_BGR2RGB),
                channels="RGB",
                use_column_width=True
            )

            time.sleep(2)  # Delay to prevent rapid looping
        except Exception as e:
            st.write(f"❌ Error: {e}")
            logging.error(f"Error during hand sign recognition: {e}")

        if cv2.waitKey(5) & 0xFF == 27:  # Press escape to break
            break

    cap.release()
    cv2.destroyAllWindows()
    stframe.empty()
    progress_bar.empty()

# Function to extract text from an audio file and translate it
def audio_to_text_from_file(audio_path, target_language, email_sender, email_address, username):
    recognizer = sr.Recognizer()
    with sr.AudioFile(audio_path) as source:
        st.write("Processing audio file...")
        audio = recognizer.record(source)
        try:
            with st.spinner("Recognizing speech..."):
                text = recognizer.recognize_google(audio)
            st.write(f"**Text extracted from audio:**")
            st.write(text)
            translated_text = translate_text(text, target_language, username)
            st.write(f"**Translated Text:**")
            st.write(translated_text)
            # Convert text to speech
            text_to_voice(translated_text, target_language)
            # Perform sentiment analysis
            sentiment_analysis(translated_text, username)
            # Log the audio file translation
            log_user_action(username, "Audio File Translation", f"Original Text: {text} | Translated Text: {translated_text}")
            return text, translated_text
        except sr.UnknownValueError:
            st.write("Speech Recognition could not understand the audio.")
            logging.error("Speech Recognition could not understand the audio.")
            # Log the failed audio file translation
            log_user_action(username, "Audio File Translation Failed", "Could not understand audio.")
            return None, None
        except sr.RequestError as e:
            st.write(f"⚠️ Could not request results; {e}")
            logging.error(f"Speech Recognition request error: {e}")
            # Log the request error
            log_user_action(username, "Audio File Translation Error", f"Request error: {e}")
            return None, None

# Function to extract text from an image with preprocessing
reader = easyocr.Reader(['en'], gpu=False)
def image_to_text(image, target_language, email_sender, email_address, username):
    try:
        # Convert PIL image → NumPy array
        img_np = np.array(image)

        # OCR extraction using EasyOCR
        result = reader.readtext(img_np, detail=0)
        text = " ".join(result)

        st.write("**Text extracted from image:**")
        st.write(text)

        # Translate text
        translated_text = translate_text(text, target_language, username)
        
        st.write("**Translated Text:**")
        st.write(translated_text)

        # Convert text to speech
        text_to_voice(translated_text, target_language)

        # Sentiment analysis
        sentiment_analysis(translated_text, username)

        # Download button
        if translated_text:
            st.download_button(
                label="Download Translated Text",
                data=translated_text,
                file_name='translated_text.txt',
                mime='text/plain'
            )

        # Logging
        log_user_action(username, "Image Text Extraction", f"Extracted Text: {text}")

        return text, translated_text

    except Exception as e:
        st.write(f"⚠️ OCR Error: {e}")
        logging.error(f"OCR Error: {e}")
        log_user_action(username, "Image Text Extraction Error", f"Error: {e}")
        return "", ""


# Function to extract text from a text file and translate it
def text_file_to_text(file, target_language, email_sender, email_address, username):
    try:
        # Read the content of the text file
        content = file.read().decode('utf-8')
        st.write(f"**Text extracted from file:**")
        st.write(content)
        translated_text = translate_text(content, target_language, username)
        st.write(f"**Translated Text:**")
        st.write(translated_text)
        # Convert text to speech
        text_to_voice(translated_text, target_language)
        # Perform sentiment analysis
        sentiment_analysis(translated_text, username)
        # Provide download option
        st.download_button(
            label="Download Translated Text",
            data=translated_text,
            file_name='translated_text.txt',
            mime='text/plain'
        )
        # Log the text file translation
        log_user_action(username, "Text File Translation", f"Original Text: {content} | Translated Text: {translated_text}")
        return content, translated_text
    except Exception as e:
        st.write(f"⚠️ Error in text file to text conversion: {e}")
        logging.error(f"Error in text file to text conversion: {e}")
        # Log the error
        log_user_action(username, "Text File Translation Error", f"Error: {e}")
        return "", ""

# Function to extract text from a Word document and translate it
def word_file_to_text(file, target_language, email_sender, email_address, username):
    try:
        st.write("Processing Word document...")
        doc = Document(file)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        text = '\n'.join(full_text)
        st.write(f"**Text extracted from Word document:**")
        st.write(text)
        translated_text = translate_text(text, target_language, username)
        st.write(f"**Translated Text:**")
        st.write(translated_text)
        # Convert text to speech
        text_to_voice(translated_text, target_language)
        # Perform sentiment analysis
        sentiment_analysis(translated_text, username)
        # Provide download option
        if text:
            st.download_button(
                label="Download Translated Text",
                data=translated_text,
                file_name='translated_text.txt',
                mime='text/plain',
                key="download_word_translation"
            )
        # Log the Word document translation
        log_user_action(username, "Word Document Translation", f"Original Text: {text} | Translated Text: {translated_text}")
        return text, translated_text
    except Exception as e:
        st.write(f"⚠️ Error in Word file to text conversion: {e}")
        logging.error(f"Error in Word file to text conversion: {e}")
        # Log the error
        log_user_action(username, "Word Document Translation Error", f"Error: {e}")
        return "", ""

# Function to extract text from a PDF file and translate it
def pdf_file_to_text(file, target_language, email_sender, email_address, username):
    try:
        st.write("Processing PDF file...")
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            extracted_text = page.extract_text()
            if extracted_text:
                text += extracted_text + "\n"
        st.write(f"**Text extracted from PDF file:**")
        st.write(text)
        translated_text = translate_text(text, target_language, username)
        st.write(f"**Translated Text:**")
        st.write(translated_text)
        # Convert text to speech
        text_to_voice(translated_text, target_language)
        # Perform sentiment analysis
        sentiment_analysis(translated_text, username)
        # Provide download option
        if text:
            st.download_button(
                label="Download Translated Text",
                data=text,
                file_name='translated_text.txt',
                mime='text/plain',
                key="download_pdf_translation"
            )
        # Log the PDF file translation
        log_user_action(username, "PDF File Translation", f"Original Text: {text} | Translated Text: {translated_text}")
        return text, translated_text
    except Exception as e:
        st.write(f"⚠️ Error in PDF file to text conversion: {e}")
        logging.error(f"Error in PDF file to text conversion: {e}")
        # Log the error
        log_user_action(username, "PDF File Translation Error", f"Error: {e}")
        return "", ""

# Function to extract text from a PowerPoint file and translate it
def ppt_file_to_text(file, target_language, email_sender, email_address, username):
    try:
        st.write("Processing PowerPoint presentation...")
        prs = Presentation(file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        text = '\n'.join(text_runs)
        st.write(f"**Text extracted from PowerPoint presentation:**")
        st.write(text)
        translated_text = translate_text(text, target_language, username)
        st.write(f"**Translated Text:**")
        st.write(translated_text)
        # Convert text to speech
        text_to_voice(translated_text, target_language)
        # Perform sentiment analysis
        sentiment_analysis(translated_text, username)
        # Provide download option
        if text:
            st.download_button(
                label="Download Translated Text",
                data=translated_text,
                file_name='translated_text.txt',
                mime='text/plain',
                key="download_ppt_translation"
            )
        # Log the PPT file translation
        log_user_action(username, "PowerPoint Presentation Translation", f"Original Text: {text} | Translated Text: {translated_text}")
        return text, translated_text
    except Exception as e:
        st.write(f"⚠️ Error in PowerPoint file to text conversion: {e}")
        logging.error(f"Error in PowerPoint file to text conversion: {e}")
        # Log the error
        log_user_action(username, "PowerPoint Presentation Translation Error", f"Error: {e}")
        return "", ""

# Function to extract text from a video file and translate it
def video_file_to_text(file, target_language, email_sender, email_address, username):
    try:
        st.write("Processing video file...")
        # Save the uploaded video file to a temporary location
        temp_dir = 'temp'
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
        temp_video_path = os.path.join(temp_dir, file.name)
        with open(temp_video_path, "wb") as f:
            f.write(file.read())

        # Extract audio from the video
        video = VideoFileClip(temp_video_path)
        temp_audio_path = os.path.join(temp_dir, f"{file.name}_audio.wav")
        video.audio.write_audiofile(temp_audio_path, codec='pcm_s16le')
        video.close()

        # Process the extracted audio
        text, translated_text = audio_to_text_from_file(
            temp_audio_path, target_language, email_sender, email_address, username
        )

        # Clean up temporary files
        os.remove(temp_video_path)
        os.remove(temp_audio_path)

        # Log the video file translation
        log_user_action(username, "Video File Translation", f"Original Text: {text} | Translated Text: {translated_text}")
        return text, translated_text
    except Exception as e:
        st.write(f"⚠️ Error in video file to text conversion: {e}")
        logging.error(f"Error in video file to text conversion: {e}")
        # Log the error
        log_user_action(username, "Video File Translation Error", f"Error: {e}")
        return "", ""

# Function to extract text from live video (webcam) and translate it
def live_video_translation(target_language, email_sender, email_address, username):
    try:
        st.write("Starting live video translation. Please allow access to your camera.")
        cap = cv2.VideoCapture(0)
        stframe = st.empty()
        progress_bar = st.progress(0)
        for percent_complete in range(100):
            time.sleep(0.01)
            progress_bar.progress(percent_complete + 1)
        while cap.isOpened():
            success, frame = cap.read()
            if not success:
                st.write("Ignoring empty camera frame.")
                continue

            # Display the frame in Streamlit
            stframe.image(
                cv2.cvtColor(frame, cv2.COLOR_BGR2RGB),
                channels="RGB",
                use_column_width=True
            )

            # For simplicity, we'll capture audio separately
            break

        cap.release()
        cv2.destroyAllWindows()
        stframe.empty()
        progress_bar.empty()

        # Capture audio after video
        recognizer = sr.Recognizer()
        with sr.Microphone() as source:
            st.write("Capturing audio from live video...")
            recognizer.adjust_for_ambient_noise(source, duration=1)
            audio = recognizer.listen(source)
            try:
                with st.spinner("Recognizing speech..."):
                    text = recognizer.recognize_google(audio)
                st.write(f"**Captured Audio Text:**")
                st.write(text)
                translated_text = translate_text(text, target_language, username)
                st.write(f"**Translated Text:**")
                st.write(translated_text)
                # Convert text to speech
                text_to_voice(translated_text, target_language)
                # Perform sentiment analysis
                sentiment_analysis(translated_text, username)
                # Option to send email
                if st.button("Send Output via Email"):
                    subject = "Live Video Translation Output"
                    body = f"Captured Audio Text: {text}\nTranslated Text: {translated_text}"
                    recipient = email_address  # Sending to the sender's email
                    send_email(subject, body, recipient, email_sender, email_address, username)
                return text, translated_text
            except sr.UnknownValueError:
                st.write("Speech Recognition could not understand the audio.")
                logging.error("Speech Recognition could not understand the audio.")
                # Log the failed live video translation
                log_user_action(username, "Live Video Translation Failed", "Could not understand audio.")
                return None, None
            except sr.RequestError as e:
                st.write(f"⚠️ Could not request results; {e}")
                logging.error(f"Speech Recognition request error: {e}")
                # Log the request error
                log_user_action(username, "Live Video Translation Error", f"Request error: {e}")
                return None, None

    except Exception as e:
        st.write(f"⚠️ Error in live video translation: {e}")
        logging.error(f"Error in live video translation: {e}")
        # Log the error
        log_user_action(username, "Live Video Translation Error", f"Error: {e}")
        return None, None

# Main entry point
if __name__ == "__main__":
    main()




